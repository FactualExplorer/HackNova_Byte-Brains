"""
Task 7: Generate Presentation Slides (7-slide PDF)
Uses python-pptx to create a professional slide deck embedding all charts and tables.
"""

import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import warnings
warnings.filterwarnings('ignore')

BASE_DIR = os.path.dirname(os.path.dirname(__file__))
CHARTS_DIR = os.path.join(BASE_DIR, 'charts')
SUBMISSIONS_DIR = os.path.join(BASE_DIR, 'submissions')

BG_DARK = RGBColor(0x0f, 0x17, 0x2a)
BG_CARD = RGBColor(0x1e, 0x29, 0x3b)
TEXT_WHITE = RGBColor(0xff, 0xff, 0xff)
TEXT_GRAY = RGBColor(0xcb, 0xd5, 0xe1)
ACCENT_CYAN = RGBColor(0x22, 0xd3, 0xee)
ACCENT_GOLD = RGBColor(0xfb, 0xbf, 0x24)
ACCENT_GREEN = RGBColor(0x22, 0xc5, 0x5e)
ACCENT_RED = RGBColor(0xef, 0x44, 0x44)
BANKING_BLUE = RGBColor(0x25, 0x63, 0xeb)
IT_GREEN = RGBColor(0x16, 0xa3, 0x4a)
PHARMA_RED = RGBColor(0xdc, 0x26, 0x26)


def set_slide_bg(slide, color=BG_DARK):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title_text, subtitle_text=None):
    """Add styled title to slide."""
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(8.6), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(4)


def add_text_box(slide, left, top, width, height, text, size=11, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_multi_text_box(slide, left, top, width, height, lines, default_size=11, default_color=TEXT_WHITE):
    """Add a text box with multiple styled lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, size, color, bold = line_info, default_size, default_color, False
        elif isinstance(line_info, tuple):
            text = line_info[0]
            size = line_info[1] if len(line_info) > 1 else default_size
            color = line_info[2] if len(line_info) > 2 else default_color
            bold = line_info[3] if len(line_info) > 3 else False
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_before = Pt(2)
    
    return tf


def add_table_to_slide(slide, df, left, top, width, height, font_size=9):
    """Add a pandas DataFrame as a styled table to the slide."""
    rows, cols = df.shape
    table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
    table = table_shape.table
    
    for j, col_name in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEXT_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x33, 0x41, 0x55)
    
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iloc[i, j])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = TEXT_GRAY
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else RGBColor(0x17, 0x22, 0x33)
    
    return table


def safe_add_picture(slide, path, left, top, width, height):
    """Add a picture to the slide if it exists."""
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    else:
        print(f"      ⚠️ Missing chart: {os.path.basename(path)}")
        return False


def create_presentation():
    """Create the 7-slide presentation."""
    print("=" * 60)
    print("GENERATING PRESENTATION SLIDES")
    print("=" * 60)
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    try:
        metrics = pd.read_csv(os.path.join(SUBMISSIONS_DIR, 'metrics_summary.csv'))
    except:
        metrics = pd.DataFrame()
    try:
        sma_signals = pd.read_csv(os.path.join(SUBMISSIONS_DIR, 'sma_signal_table.csv'))
    except:
        sma_signals = pd.DataFrame()
    try:
        portfolio = pd.read_csv(os.path.join(SUBMISSIONS_DIR, 'portfolio_comparison.csv'))
    except:
        portfolio = pd.DataFrame()
    try:
        chaos = pd.read_csv(os.path.join(SUBMISSIONS_DIR, 'chaos_round_stress_test.csv'))
    except:
        chaos = pd.DataFrame()
    try:
        ml_comp = pd.read_csv(os.path.join(SUBMISSIONS_DIR, 'ml_comparison_table.csv'))
    except:
        ml_comp = pd.DataFrame()
    try:
        sector_breakdown = pd.read_csv(os.path.join(SUBMISSIONS_DIR, 'sector_breakdown.csv'))
    except:
        sector_breakdown = pd.DataFrame()
    
    print("  📄 Slide 1: Data Overview...")
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)
    add_title(slide1, "📊 Slide 1 — Data Acquisition & Quality",
              "Source, Period, Quality Issues & Resolutions")
    
    data_lines = [
        ("DATA SOURCE & PERIOD", 14, ACCENT_CYAN, True),
        ("• Provider: Yahoo Finance (yfinance Python library)", 14, TEXT_GRAY, False),
        ("• Period: 01 Jan 2023 – 31 Dec 2024 (2 full calendar years)", 14, TEXT_GRAY, False),
        ("• Frequency: Daily OHLCV (Adjusted Close for all return calculations)", 14, TEXT_GRAY, False),
        ("• Universe: 15 NSE-listed stocks across 3 sectors", 14, TEXT_GRAY, False),
        ("    Banking (5): HDFCBANK, ICICIBANK, SBIN, KOTAKBANK, AXISBANK", 12, TEXT_GRAY, False),
        ("    IT (5): TCS, INFY, WIPRO, HCLTECH, TECHM", 12, TEXT_GRAY, False),
        ("    Pharma (5): SUNPHARMA, DRREDDY, CIPLA, DIVISLAB, APOLLOHOSP", 12, TEXT_GRAY, False),
        ("• Benchmark: Nifty 50 Index (^NSEI)", 14, TEXT_GRAY, False),
        ("", 8, TEXT_GRAY, False),
        ("DATA QUALITY ACTIONS", 14, ACCENT_GOLD, True),
        ("• Missing values → Forward-filled, then back-filled", 13, TEXT_GRAY, False),
        ("• Indian market holidays → Expected gaps, no additional treatment", 13, TEXT_GRAY, False),
        ("• Zero/negative prices → Replaced with forward-fill", 13, TEXT_GRAY, False),
        ("• Extreme returns (>25% daily) → Flagged but retained (legitimate events)", 13, TEXT_GRAY, False),
        ("• Validated: No stock has >5 consecutive missing trading days", 13, TEXT_GRAY, False),
        ("• All timestamps timezone-normalized (IST assumed)", 13, TEXT_GRAY, False),
    ]
    add_multi_text_box(slide1, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5), data_lines)
    
    print("  📄 Slide 2: Risk-Return Map...")
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_title(slide2, "📈 Slide 2 — Risk-Return Map",
              "Annualized Volatility vs Annualized Return — Sector Clusters Explained")
    
    safe_add_picture(slide2, os.path.join(CHARTS_DIR, 'risk_return_scatter.png'),
                     Inches(0.3), Inches(1.2), Inches(8.5), Inches(5.8))
    
    safe_add_picture(slide2, os.path.join(CHARTS_DIR, 'correlation_heatmap.png'),
                     Inches(9), Inches(1.2), Inches(4), Inches(3.5))
    
    add_text_box(slide2, Inches(9), Inches(4.9), Inches(4), Inches(0.3),
                 "Daily Returns Correlation Heatmap", size=10, color=TEXT_GRAY, bold=True)
    
    insight_lines = [
        ("KEY OBSERVATIONS", 11, ACCENT_CYAN, True),
        ("• Banking cluster: Moderate vol, correlated", 10, TEXT_GRAY, False),
        ("• IT cluster: Varies widely in risk-return", 10, TEXT_GRAY, False),
        ("• Pharma: Defensive, lower Beta", 10, TEXT_GRAY, False),
        ("• Risk-free rate line at 6.5% p.a.", 10, TEXT_GRAY, False),
    ]
    add_multi_text_box(slide2, Inches(9), Inches(5.3), Inches(4), Inches(2), insight_lines)
    
    print("  📄 Slide 3: Top & Bottom Stocks...")
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_title(slide3, "🏆 Slide 3 — Top & Bottom Stocks by Sharpe Ratio",
              "Best 3 and Worst 3 performers — What separates winners from laggards")
    
    if not metrics.empty:
        display_cols = ['Ticker', 'Sector', 'Annualized Return (%)',
                        'Annualized Volatility (%)', 'Sharpe Ratio', 'Beta vs Nifty 50', 'Maximum Drawdown (%)']
        sorted_m = metrics.sort_values('Sharpe Ratio', ascending=False)
        top3 = sorted_m.head(3)[display_cols]
        bot3 = sorted_m.tail(3)[display_cols]
        
        add_text_box(slide3, Inches(0.7), Inches(1.4), Inches(4), Inches(0.4),
                     "🟢 TOP 3 (Highest Sharpe Ratio)", size=16, color=ACCENT_GREEN, bold=True)
        add_table_to_slide(slide3, top3.reset_index(drop=True),
                          Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.5), font_size=11)
        
        add_text_box(slide3, Inches(0.7), Inches(3.8), Inches(4), Inches(0.4),
                     "🔴 BOTTOM 3 (Lowest Sharpe Ratio)", size=16, color=ACCENT_RED, bold=True)
        add_table_to_slide(slide3, bot3.reset_index(drop=True),
                          Inches(0.5), Inches(4.3), Inches(12.3), Inches(1.5), font_size=11)
        
        if not sector_breakdown.empty:
            add_text_box(slide3, Inches(0.7), Inches(6.1), Inches(4), Inches(0.3),
                        "📊 Sector Averages", size=13, color=ACCENT_CYAN, bold=True)
            add_table_to_slide(slide3, sector_breakdown,
                              Inches(0.5), Inches(6.4), Inches(8), Inches(0.9), font_size=10)
    
    print("  📄 Slide 4: SMA Signals...")
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_title(slide4, "📉 Slide 4 — SMA Technical Signals",
              "50-Day & 200-Day SMA — Golden Cross / Death Cross as of 31 Dec 2024")
    
    if not sma_signals.empty:
        add_table_to_slide(slide4, sma_signals,
                          Inches(0.3), Inches(1.4), Inches(7.5), Inches(4.5), font_size=10)
    
    sma_chart = os.path.join(CHARTS_DIR, 'sma_banking.png')
    if safe_add_picture(slide4, sma_chart, Inches(8), Inches(1.4), Inches(5), Inches(2.8)):
        add_text_box(slide4, Inches(8), Inches(4.3), Inches(5), Inches(0.3),
                     "HDFCBANK — SMA Overlay (Banking)", size=9, color=TEXT_GRAY, bold=True)
    
    sma_chart2 = os.path.join(CHARTS_DIR, 'sma_it.png')
    if safe_add_picture(slide4, sma_chart2, Inches(8), Inches(4.6), Inches(5), Inches(2.5)):
        add_text_box(slide4, Inches(8), Inches(7.0), Inches(5), Inches(0.3),
                     "TCS — SMA Overlay (IT)", size=9, color=TEXT_GRAY, bold=True)
    
    print("  📄 Slide 5: Portfolio Recommendation + Chaos Round...")
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_title(slide5, "💼 Slide 5 — Portfolio Recommendation + Chaos Round",
              "Portfolio A (Equal) vs Portfolio B (Optimized) + Nifty -10% Stress Test")
    
    if not portfolio.empty:
        add_text_box(slide5, Inches(0.5), Inches(1.3), Inches(4), Inches(0.3),
                    "Portfolio Comparison", size=13, color=ACCENT_CYAN, bold=True)
        add_table_to_slide(slide5, portfolio,
                          Inches(0.5), Inches(1.7), Inches(5.5), Inches(2.3), font_size=10)
    
    safe_add_picture(slide5, os.path.join(CHARTS_DIR, 'portfolio_comparison.png'),
                     Inches(6.3), Inches(1.2), Inches(6.8), Inches(2.8))
    
    add_text_box(slide5, Inches(0.5), Inches(4.2), Inches(5), Inches(0.3),
                "⚡ Chaos Round: Nifty -10% Stress Test", size=13, color=ACCENT_GOLD, bold=True)
    
    chaos_chart = os.path.join(CHARTS_DIR, 'chaos_stress_test.png')
    if safe_add_picture(slide5, chaos_chart, Inches(0.3), Inches(4.6), Inches(6.2), Inches(2.7)):
        pass
    
    # Chaos data table (compact)
    if not chaos.empty:
        chaos_key = chaos[chaos['Ticker'].isin(['PORTFOLIO A', 'PORTFOLIO B', 'MOST EXPOSED', 'SAFEST REFUGE'])]
        if not chaos_key.empty:
            add_table_to_slide(slide5, chaos_key[['Ticker', 'Sector', 'Beta', 'Expected Loss (%)']].reset_index(drop=True),
                              Inches(6.5), Inches(4.6), Inches(6.3), Inches(1.8), font_size=10)
    
    add_text_box(slide5, Inches(6.5), Inches(6.6), Inches(6.3), Inches(0.5),
                 "⚠️ Portfolio justification (200 words) must be presented live",
                 size=11, color=ACCENT_GOLD)
    
    print("  📄 Slide 6: ML Model...")
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)
    add_title(slide6, "🤖 Slide 6 — ML Model Comparison",
              "Random Forest vs Gradient Boosting — Predicting Stock Outperformance vs Nifty 50 (20-day forward)")
    
    # ML comparison table
    if not ml_comp.empty:
        add_text_box(slide6, Inches(0.5), Inches(1.3), Inches(4), Inches(0.3),
                    "Model Performance (Test: 2024)", size=12, color=ACCENT_CYAN, bold=True)
        add_table_to_slide(slide6, ml_comp,
                          Inches(0.5), Inches(1.65), Inches(7), Inches(1.0), font_size=11)
    
    safe_add_picture(slide6, os.path.join(CHARTS_DIR, 'ml_feature_importance.png'),
                     Inches(0.3), Inches(2.9), Inches(6.5), Inches(2.5))
    
    safe_add_picture(slide6, os.path.join(CHARTS_DIR, 'ml_confusion_matrix.png'),
                     Inches(6.8), Inches(1.3), Inches(6.3), Inches(2.8))
    
    safe_add_picture(slide6, os.path.join(CHARTS_DIR, 'ml_roc_curve.png'),
                     Inches(6.8), Inches(4.2), Inches(6.3), Inches(3))
    
    add_text_box(slide6, Inches(0.5), Inches(5.6), Inches(6.5), Inches(1.5),
                 "ML Design:\n"
                 "• Target: Does stock outperform Nifty 50 over next 20 days?\n"
                 "• Features: 18 technical indicators (returns, volatility, SMA ratios,\n"
                 "  RSI, MACD, Beta, volume)\n"
                 "• Temporal split: Train on 2023, Test on 2024 (no data leakage)\n"
                 "• Class balancing: balanced weights (RF) / scale_pos_weight (XGB)",
                 size=11, color=TEXT_GRAY)
    
    add_text_box(slide6, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                 "⚠️ Top 3 features and financial interpretation must be explained live by YOU (AI policy)",
                 size=11, color=ACCENT_GOLD)
    
    print("  📄 Slide 7: n8n Demo...")
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7)
    add_title(slide7, "⚙️ Slide 7 — n8n Automation Demo",
              "Price Alert Workflow — Live Execution")
    
    n8n_diagram = os.path.join(CHARTS_DIR, 'n8n_workflow_diagram.png')
    if not safe_add_picture(slide7, n8n_diagram, Inches(0.3), Inches(1.3), Inches(8), Inches(4.5)):
        # Fallback: try n8n directory
        n8n_diagram2 = os.path.join(BASE_DIR, 'n8n', 'n8n_workflow_diagram.png')
        safe_add_picture(slide7, n8n_diagram2, Inches(0.3), Inches(1.3), Inches(8), Inches(4.5))
    
    add_multi_text_box(slide7, Inches(8.5), Inches(1.3), Inches(4.5), Inches(5.5), n8n_lines)
    
    add_text_box(slide7, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
                 "Workflow JSON → n8n/price_alert_workflow.json  |  "
                 "Runs Mon-Fri 4 PM IST  |  Uses Yahoo Finance API  |  "
                 "Sends alert via email to fund manager",
                 size=11, color=TEXT_GRAY)
    
    pptx_path = os.path.join(SUBMISSIONS_DIR, 'presentation.pptx')
    prs.save(pptx_path)
    print(f"\n💾 Presentation saved: {pptx_path}")
    print("📌 Note: Export the .pptx to PDF for final submission")
    print("📌 Note: Export the .pptx to PDF for final submission")
    
    return pptx_path


if __name__ == '__main__':
    pptx_path = create_presentation()
    
    print("\n" + "=" * 60)
    print("✅ PRESENTATION COMPLETE")
    print("=" * 60)
    print(f"\nDeliverables:")
    print(f"  📄 {pptx_path}")
    print(f"\n📌 Remember to export to PDF before submission")
    print(f"⚠️  Review all slides and add your own financial interpretations")
